"""
Извлечение исходного текста из документов и URL для режима
source_type != TOPIC (структуризация чужого материала, см. generation/llm.py).

Ответственность:
- Достать текстовый слой из .pdf/.docx/.pptx/.txt
- Достать основной текст статьи из списка URL
- Всегда вернуть текст, обрезанный до MAX_CHARS — дальше в промпт идёт как есть
"""

import asyncio
import logging

import chardet
import httpx
import trafilatura
from pdfplumber import open as pdfplumber_open
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation

logger = logging.getLogger(__name__)

MAX_CHARS = 15000
MAX_PDF_PAGES = 40
URL_TIMEOUT_SECONDS = 10
MAX_CONCURRENT_URLS = 3

SUPPORTED_MIME_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "text/plain": "txt",
}


def _truncate(text: str, source: str) -> str:
    if len(text) > MAX_CHARS:
        logger.warning(
            "Content truncated to %d chars",
            MAX_CHARS,
            extra={"source": source, "original_length": len(text)},
        )
        return text[:MAX_CHARS]
    return text


def _extract_pdf(file_bytes: bytes) -> str:
    import io

    chunks: list[str] = []
    with pdfplumber_open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages[:MAX_PDF_PAGES]:
            text = page.extract_text()
            if text:
                chunks.append(text)
    return "\n\n".join(chunks)


def _extract_docx(file_bytes: bytes) -> str:
    import io

    doc = DocxDocument(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _extract_pptx(file_bytes: bytes) -> str:
    import io

    prs = PptxPresentation(io.BytesIO(file_bytes))
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                chunks.append(notes)
    return "\n".join(chunks)


def _extract_txt(file_bytes: bytes) -> str:
    detected = chardet.detect(file_bytes)
    encoding = detected.get("encoding") or "utf-8"
    return file_bytes.decode(encoding, errors="replace")


_EXTRACTORS = {
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "txt": _extract_txt,
}


async def extract_from_document(file_bytes: bytes, mime_type: str) -> str:
    """
    Извлекает текст из файла по mime_type.

    Raises:
        ValueError: неподдерживаемый mime_type
    """
    fmt = SUPPORTED_MIME_TYPES.get(mime_type)
    if fmt is None:
        raise ValueError(f"Unsupported mime_type: {mime_type}")

    extractor = _EXTRACTORS[fmt]
    text = await asyncio.to_thread(extractor, file_bytes)
    return _truncate(text.strip(), source=f"document:{fmt}")


async def _fetch_one_url(client: httpx.AsyncClient, url: str) -> tuple[str, str | None]:
    """Возвращает (url, текст|None). None означает ошибку — url пойдёт в error_urls."""
    try:
        response = await client.get(url, timeout=URL_TIMEOUT_SECONDS, follow_redirects=True)
        response.raise_for_status()
        extracted = trafilatura.extract(response.text)
        if not extracted:
            logger.warning("No extractable content", extra={"url": url})
            return url, None
        return url, extracted
    except Exception as e:
        logger.warning("Failed to fetch URL", extra={"url": url, "error": str(e)})
        return url, None


async def extract_from_url(urls: list[str]) -> tuple[str, list[str]]:
    """
    Достаёт основной текст статьи для каждого URL (максимум MAX_CONCURRENT_URLS
    одновременно). Одна битая ссылка не валит остальные.

    Returns:
        (объединённый_текст, список_ссылок_с_ошибкой)
    """
    urls = urls[:MAX_CONCURRENT_URLS]
    semaphore = asyncio.Semaphore(MAX_CONCURRENT_URLS)

    async def _bounded_fetch(client: httpx.AsyncClient, url: str):
        async with semaphore:
            return await _fetch_one_url(client, url)

    async with httpx.AsyncClient(headers={"User-Agent": "Mozilla/5.0 (compatible; FibonacciAI/1.0)"}) as client:
        results = await asyncio.gather(
            *(_bounded_fetch(client, url) for url in urls),
            return_exceptions=True,
        )

    texts: list[str] = []
    error_urls: list[str] = []
    for i, result in enumerate(results):
        if isinstance(result, Exception):
            error_urls.append(urls[i])
            continue
        url, text = result
        if text:
            texts.append(f"--- источник: {url} ---\n{text}")
        else:
            error_urls.append(url)

    combined = "\n\n".join(texts)
    return _truncate(combined, source="url"), error_urls
